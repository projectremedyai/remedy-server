"""Stage 3: Content Extraction — OCR and native parsing for all document types.

Routes each document to the appropriate extraction strategy:
  - PDF  -> vision OCR layout parsing
  - DOCX -> python-docx native parser (vision OCR fallback for image-heavy docs)
  - PPTX -> python-pptx native parser (vision OCR fallback)
  - XLSX -> openpyxl native parser (sheet-by-sheet markdown tables)

All extractors emit clean Markdown preserving headings, lists, tables,
emphasis, and image placeholders.
"""

from __future__ import annotations

import logging
import re
from pathlib import Path
from typing import Any

from project_remedy.config import PipelineConfig
from project_remedy.database import DatabaseManager
from project_remedy.image_extractor import extract_pdf_images
from project_remedy.models import DocumentJob, ExtractedImage, FileType, JobStatus
from project_remedy.ollama_client import OllamaClient, OllamaClientError

logger = logging.getLogger(__name__)

# Threshold: if a DOCX/PPTX has this many embedded images relative to
# text paragraphs, we fall back to vision OCR for more accurate extraction.
_IMAGE_HEAVY_RATIO = 0.5


class ExtractionError(Exception):
    """Raised when content extraction fails for a document."""


class ContentExtractor:
    """Orchestrates content extraction across all supported file types.

    Parameters
    ----------
    config:
        Pipeline configuration.
    ollama:
        Initialised Ollama API client.
    db:
        Database manager for persisting job state.
    """

    def __init__(
        self,
        config: PipelineConfig,
        ollama: OllamaClient,
        db: DatabaseManager,
    ) -> None:
        self._config = config
        self._ollama = ollama
        self._db = db

    # ------------------------------------------------------------------
    # Public entry point
    # ------------------------------------------------------------------

    async def extract(self, job: DocumentJob) -> str:
        """Extract content from the downloaded document.

        Updates the job status through EXTRACTING -> EXTRACTED (or FAILED)
        and stores the result in ``job.ocr_markdown``.

        Returns
        -------
        str
            Extracted Markdown content.
        """
        job.status = JobStatus.EXTRACTING
        await self._db.update_job(job)
        logger.info("Extracting content for job %s (%s)", job.id, job.file_type)

        try:
            file_path = Path(job.local_path)
            if not file_path.exists():
                raise ExtractionError(f"Local file not found: {file_path}")

            file_type = job.file_type
            if file_type is None:
                raise ExtractionError("Job has no file_type set.")

            if file_type == FileType.PDF:
                markdown = await self._extract_pdf(file_path, job)
            elif file_type in (FileType.DOCX, FileType.DOC):
                markdown = await self._extract_docx(file_path)
            elif file_type in (FileType.PPTX, FileType.PPT):
                markdown = await self._extract_pptx(file_path)
            elif file_type in (FileType.XLSX, FileType.XLS):
                markdown = await self._extract_xlsx(file_path)
            else:
                raise ExtractionError(f"Unsupported file type: {file_type}")

            if not markdown.strip():
                raise ExtractionError("Extraction produced empty content.")

            job.ocr_markdown = markdown
            job.status = JobStatus.EXTRACTED
            job.error_message = ""
            await self._db.update_job(job)
            logger.info(
                "Extraction complete for job %s — %d chars of markdown",
                job.id,
                len(markdown),
            )
            return markdown

        except Exception as exc:
            error_msg = f"Extraction failed: {exc}"
            logger.error("Job %s: %s", job.id, error_msg)
            job.status = JobStatus.FAILED
            job.error_message = error_msg
            await self._db.update_job(job)
            raise ExtractionError(error_msg) from exc

    # ------------------------------------------------------------------
    # PDF extraction
    # ------------------------------------------------------------------

    async def _extract_pdf(self, file_path: Path, job: DocumentJob) -> str:
        """Extract PDF content via vision OCR with image extraction."""
        logger.debug("PDF extraction via vision OCR: %s", file_path.name)

        # Extract embedded images from the PDF.
        images_dir = self._config.output.output_dir / "images" / job.id[:12]
        extracted = extract_pdf_images(file_path, images_dir)
        job.set_extracted_images(extracted)

        # Run OCR.
        markdown = await self._ollama.ocr(file_path=file_path)

        # Replace hallucinated image URLs with actual extracted filenames.
        if extracted:
            markdown = _sanitize_image_refs(markdown, extracted)

        return markdown

    # ------------------------------------------------------------------
    # DOCX extraction
    # ------------------------------------------------------------------

    async def _extract_docx(self, file_path: Path) -> str:
        """Extract DOCX content using python-docx, with vision OCR fallback."""
        try:
            from docx import Document as DocxDocument  # type: ignore[import-untyped]
        except ImportError:
            logger.warning(
                "python-docx not installed — falling back to vision OCR for %s",
                file_path.name,
            )
            return await self._ollama.ocr(file_path=file_path)

        logger.debug("DOCX native extraction: %s", file_path.name)
        doc = DocxDocument(str(file_path))

        # Count images to decide if we should fall back to OCR.
        image_count = sum(
            1
            for rel in doc.part.rels.values()
            if "image" in rel.reltype
        )
        paragraph_count = len(doc.paragraphs)
        if paragraph_count > 0 and image_count / paragraph_count > _IMAGE_HEAVY_RATIO:
            logger.info(
                "DOCX %s is image-heavy (%d images / %d paragraphs) "
                "— falling back to vision OCR",
                file_path.name,
                image_count,
                paragraph_count,
            )
            return await self._ollama.ocr(file_path=file_path)

        parts: list[str] = []
        image_index = 0

        for element in doc.element.body:
            tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

            if tag == "p":
                text = self._docx_paragraph_to_md(element, doc)
                if text:
                    parts.append(text)

            elif tag == "tbl":
                table_md = self._docx_table_to_md(element, doc)
                if table_md:
                    parts.append(table_md)

        # Append image placeholders.
        for i in range(image_count):
            parts.append(f"\n![Image {i + 1}](image_{i + 1}.png)\n")

        return "\n\n".join(parts)

    def _docx_paragraph_to_md(self, para_element: Any, doc: Any) -> str:
        """Convert a DOCX paragraph XML element to Markdown."""
        from docx.text.paragraph import Paragraph  # type: ignore[import-untyped]

        para = Paragraph(para_element, doc)
        text = ""

        # Gather run-level formatting.
        for run in para.runs:
            run_text = run.text or ""
            if not run_text:
                continue
            if run.bold:
                run_text = f"**{run_text}**"
            if run.italic:
                run_text = f"*{run_text}*"
            text += run_text

        text = text.strip()
        if not text:
            return ""

        # Determine heading level from style name.
        style_name = (para.style.name or "").lower() if para.style else ""
        if style_name.startswith("heading"):
            # Extract heading level number.
            match = re.search(r"\d+", style_name)
            level = int(match.group()) if match else 1
            level = min(level, 6)
            return f"{'#' * level} {text}"

        # Detect list styles.
        if style_name.startswith("list bullet") or style_name.startswith("list"):
            return f"- {text}"
        if style_name.startswith("list number"):
            return f"1. {text}"

        return text

    def _docx_table_to_md(self, tbl_element: Any, doc: Any) -> str:
        """Convert a DOCX table XML element to a Markdown table."""
        from docx.table import Table  # type: ignore[import-untyped]

        table = Table(tbl_element, doc)
        rows = table.rows
        if not rows:
            return ""

        md_rows: list[str] = []
        for i, row in enumerate(rows):
            cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
            md_rows.append("| " + " | ".join(cells) + " |")
            if i == 0:
                md_rows.append("| " + " | ".join("---" for _ in cells) + " |")

        return "\n".join(md_rows)

    # ------------------------------------------------------------------
    # PPTX extraction
    # ------------------------------------------------------------------

    async def _extract_pptx(self, file_path: Path) -> str:
        """Extract PPTX content slide-by-slide using python-pptx."""
        try:
            from pptx import Presentation  # type: ignore[import-untyped]
        except ImportError:
            logger.warning(
                "python-pptx not installed — falling back to vision OCR for %s",
                file_path.name,
            )
            return await self._ollama.ocr(file_path=file_path)

        logger.debug("PPTX native extraction: %s", file_path.name)
        prs = Presentation(str(file_path))

        # Check if image-heavy.
        total_shapes = 0
        image_shapes = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                total_shapes += 1
                if shape.shape_type and shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    image_shapes += 1

        if total_shapes > 0 and image_shapes / total_shapes > _IMAGE_HEAVY_RATIO:
            logger.info(
                "PPTX %s is image-heavy (%d/%d shapes are images) "
                "— falling back to vision OCR",
                file_path.name,
                image_shapes,
                total_shapes,
            )
            return await self._ollama.ocr(file_path=file_path)

        parts: list[str] = []
        image_index = 0

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_parts: list[str] = [f"## Slide {slide_num}"]

            for shape in slide.shapes:
                # Title shape.
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = self._pptx_paragraph_to_md(para, shape)
                        if text:
                            slide_parts.append(text)

                # Table shape.
                if shape.has_table:
                    table_md = self._pptx_table_to_md(shape.table)
                    if table_md:
                        slide_parts.append(table_md)

                # Image shape.
                if shape.shape_type and shape.shape_type == 13:
                    image_index += 1
                    slide_parts.append(
                        f"![Slide {slide_num} Image {image_index}]"
                        f"(slide_{slide_num}_image_{image_index}.png)"
                    )

            # Slide notes.
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_parts.append(f"\n> **Speaker Notes:** {notes_text}")

            parts.append("\n\n".join(slide_parts))

        return "\n\n---\n\n".join(parts)

    def _pptx_paragraph_to_md(self, para: Any, shape: Any) -> str:
        """Convert a python-pptx paragraph to Markdown."""
        text = ""
        for run in para.runs:
            run_text = run.text or ""
            if not run_text:
                continue
            if run.font.bold:
                run_text = f"**{run_text}**"
            if run.font.italic:
                run_text = f"*{run_text}*"
            text += run_text

        text = text.strip()
        if not text:
            return ""

        # Check if this is the title placeholder.
        is_title = False
        if hasattr(shape, "placeholder_format") and shape.placeholder_format:
            ph_idx = shape.placeholder_format.idx
            if ph_idx in (0, 1):  # Title or Center Title
                is_title = True

        if is_title:
            return f"### {text}"

        # Bullet level.
        level = para.level if hasattr(para, "level") and para.level else 0
        if level > 0:
            indent = "  " * (level - 1)
            return f"{indent}- {text}"

        return text

    def _pptx_table_to_md(self, table: Any) -> str:
        """Convert a python-pptx table to a Markdown table."""
        rows = table.rows
        if not rows:
            return ""

        md_rows: list[str] = []
        for i, row in enumerate(rows):
            cells = [
                cell.text.strip().replace("|", "\\|") for cell in row.cells
            ]
            md_rows.append("| " + " | ".join(cells) + " |")
            if i == 0:
                md_rows.append("| " + " | ".join("---" for _ in cells) + " |")

        return "\n".join(md_rows)

    # ------------------------------------------------------------------
    # XLSX extraction
    # ------------------------------------------------------------------

    async def _extract_xlsx(self, file_path: Path) -> str:
        """Extract XLSX content sheet-by-sheet as Markdown tables."""
        try:
            from openpyxl import load_workbook  # type: ignore[import-untyped]
        except ImportError:
            logger.warning(
                "openpyxl not installed — falling back to vision OCR for %s",
                file_path.name,
            )
            return await self._ollama.ocr(file_path=file_path)

        logger.debug("XLSX native extraction: %s", file_path.name)
        wb = load_workbook(str(file_path), read_only=True, data_only=True)
        parts: list[str] = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_parts: list[str] = [f"## {sheet_name}"]

            rows_data: list[list[str]] = []
            for row in ws.iter_rows(values_only=True):
                cells = [
                    str(cell).strip() if cell is not None else ""
                    for cell in row
                ]
                # Skip completely empty rows.
                if any(cells):
                    rows_data.append(cells)

            if not rows_data:
                sheet_parts.append("*Empty sheet*")
                parts.append("\n\n".join(sheet_parts))
                continue

            # Normalise column count (pad shorter rows).
            max_cols = max(len(r) for r in rows_data)
            for row in rows_data:
                while len(row) < max_cols:
                    row.append("")

            # Build Markdown table.
            for i, row in enumerate(rows_data):
                escaped = [c.replace("|", "\\|") for c in row]
                sheet_parts.append("| " + " | ".join(escaped) + " |")
                if i == 0:
                    sheet_parts.append(
                        "| " + " | ".join("---" for _ in escaped) + " |"
                    )

            parts.append("\n".join(sheet_parts))

        wb.close()
        return "\n\n---\n\n".join(parts)


# ---------------------------------------------------------------------------
# Post-processor: replace hallucinated image refs with real extracted images
# ---------------------------------------------------------------------------

_MD_IMAGE_RE = re.compile(r"!\[([^\]]*)\]\(([^)]+)\)")

_PAGE_MARKER_RE = re.compile(r"<!--\s*Page\s+(\d+)\s*-->")


def _sanitize_image_refs(
    markdown: str, extracted: list[ExtractedImage]
) -> str:
    """Replace hallucinated or placeholder image refs with real filenames.

    The OCR model often invents URLs like ``https://i.imgur.com/12345.png``
    or uses placeholder tokens.  This function splits the markdown by page
    markers, then for each page assigns extracted images in order of
    appearance.  Any extracted images that the OCR failed to reference at
    all are appended to the end of their respective page section.
    """
    if not extracted:
        return markdown

    known_filenames = {img.filename for img in extracted}

    # Build a mapping: page_number -> list of extracted images (in order).
    by_page: dict[int, list[ExtractedImage]] = {}
    for img in extracted:
        by_page.setdefault(img.page_number, []).append(img)

    # Track which images get consumed during replacement.
    used_filenames: set[str] = set()

    # Split the markdown into per-page sections.
    sections = _PAGE_MARKER_RE.split(markdown)

    # sections alternates between content and page numbers:
    # [pre-content, "1", page-1-content, "2", page-2-content, ...]
    rebuilt: list[str] = []
    current_page = 1  # default if no page markers at all

    for i, section in enumerate(sections):
        if i % 2 == 1:
            # This is a page number string from the regex capture group.
            current_page = int(section)
            rebuilt.append(f"<!-- Page {current_page} -->")
            continue

        page_images = list(by_page.get(current_page, []))
        img_idx = 0

        def _replace_image(m: re.Match) -> str:
            nonlocal img_idx
            alt_text = m.group(1)
            original_path = m.group(2)

            # If it already points to a valid extracted filename, keep it.
            if original_path in known_filenames:
                used_filenames.add(original_path)
                return m.group(0)

            # Assign the next extracted image for this page.
            if img_idx < len(page_images):
                img = page_images[img_idx]
                img_idx += 1
                used_filenames.add(img.filename)
                return f"![{alt_text}]({img.filename})"

            # No more extracted images — mark as missing.
            return f"![{alt_text}](MISSING_IMAGE)"

        section = _MD_IMAGE_RE.sub(_replace_image, section)

        # Append any extracted images that the OCR missed entirely.
        # Skip empty/whitespace-only sections (e.g. pre-content before first page marker).
        unused_on_page = [
            img for img in page_images if img.filename not in used_filenames
        ]
        if unused_on_page and section.strip():
            appendix = "\n"
            for img in unused_on_page:
                used_filenames.add(img.filename)
                appendix += f"\n![Image from page {current_page}]({img.filename})\n"
            section = section.rstrip() + appendix

        rebuilt.append(section)

    return "".join(rebuilt)
