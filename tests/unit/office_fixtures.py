"""In-memory OOXML fixture builders (NFR5: no checked-in binary blobs).

The anchored-image conversion is schema-loose (it renames ``wp:inline`` to
``wp:anchor`` without adding the positioning children Word itself would
require). That is sufficient here: these fixtures only need to be readable by
python-docx and ``xml.etree.ElementTree``, not openable in Word.
"""

from __future__ import annotations

import base64
import shutil
import zipfile
from io import BytesIO
from pathlib import Path
from typing import Sequence

# 1x1 transparent PNG
TINY_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJ"
    "AAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="
)

_LEGACY_OLE2_MAGIC = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"


def make_docx(
    path: Path,
    *,
    title: str = "",
    language: str = "",
    headings: Sequence[tuple[str, int]] = (),
    outline_paragraphs: Sequence[tuple[str, int]] = (),
    body_paragraphs: Sequence[str] = (),
    body_first: bool = False,
    tables: int = 0,
    mark_table_headers: bool = True,
    merge_header_cells: bool = False,
    inline_images: int = 0,
    image_alt: str | None = "A sample image",
    anchored_images: bool = False,
    manual_bullets: Sequence[str] = (),
    real_list_items: Sequence[str] = (),
    hyperlinks: Sequence[tuple[str, str]] = (),
    color_paragraph: str = "",
) -> Path:
    """Build a .docx exercising exactly the features the caller asks for.

    ``headings`` is a list of ``(text, level)``; level 0 applies the Title
    style, level N >= 1 applies "Heading N". ``outline_paragraphs`` is a list
    of ``(text, outline_val)``; each adds a Normal-styled paragraph (no
    heading style) carrying an explicit ``w:outlineLvl`` element, so the
    paragraph is only structurally a heading via outline level. ``body_first``
    puts one body paragraph before the first heading (for OOXML-DOCX-2.3 Fail
    fixtures). ``anchored_images=True`` converts every image to a floating
    ``wp:anchor``. ``image_alt=None`` leaves images with no descr/title.
    """
    from docx import Document
    from docx.opc.constants import RELATIONSHIP_TYPE as RT
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.shared import RGBColor

    doc = Document()
    props = doc.core_properties
    if title:
        props.title = title
    if language:
        props.language = language

    body_iter = list(body_paragraphs)
    if body_first and body_iter:
        doc.add_paragraph(body_iter.pop(0))

    for text, level in headings:
        style = "Title" if level == 0 else f"Heading {level}"
        doc.add_paragraph(text, style=style)

    for text, outline_val in outline_paragraphs:
        para = doc.add_paragraph(text)  # default Normal style, no heading style
        p_pr = para._p.get_or_add_pPr()
        outline = OxmlElement("w:outlineLvl")
        outline.set(qn("w:val"), str(outline_val))
        p_pr.append(outline)

    for text in body_iter:
        doc.add_paragraph(text)

    for text in manual_bullets:
        doc.add_paragraph(text)  # visual bullet chars, deliberately no numPr

    for text in real_list_items:
        para = doc.add_paragraph(text, style="List Bullet")
        p_pr = para._p.get_or_add_pPr()
        num_pr = OxmlElement("w:numPr")
        ilvl = OxmlElement("w:ilvl")
        ilvl.set(qn("w:val"), "0")
        num_id = OxmlElement("w:numId")
        num_id.set(qn("w:val"), "1")
        num_pr.append(ilvl)
        num_pr.append(num_id)
        p_pr.append(num_pr)

    for display_text, url in hyperlinks:
        para = doc.add_paragraph()
        r_id = para.part.relate_to(url, RT.HYPERLINK, is_external=True)
        hyperlink = OxmlElement("w:hyperlink")
        hyperlink.set(qn("r:id"), r_id)
        run = OxmlElement("w:r")
        t = OxmlElement("w:t")
        t.text = display_text
        run.append(t)
        hyperlink.append(run)
        para._p.append(hyperlink)

    if color_paragraph:
        para = doc.add_paragraph()
        run = para.add_run(color_paragraph)
        run.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)

    for _ in range(tables):
        table = doc.add_table(rows=2, cols=2)
        table.rows[0].cells[0].text = "Header A"
        table.rows[0].cells[1].text = "Header B"
        table.rows[1].cells[0].text = "data 1"
        table.rows[1].cells[1].text = "data 2"
        if merge_header_cells:
            table.rows[0].cells[0].merge(table.rows[0].cells[1])
        if mark_table_headers:
            tr_pr = table.rows[0]._tr.get_or_add_trPr()
            if tr_pr.find(qn("w:tblHeader")) is None:
                tr_pr.append(OxmlElement("w:tblHeader"))

    for _ in range(inline_images):
        doc.add_picture(BytesIO(TINY_PNG))
        if image_alt is not None:
            doc_pr = doc.inline_shapes[-1]._inline.docPr
            doc_pr.set("descr", image_alt)

    path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(path))

    if anchored_images:
        _convert_inline_images_to_anchored(path)
    return path


def _convert_inline_images_to_anchored(path: Path) -> None:
    """Rewrite word/document.xml so every wp:inline becomes wp:anchor."""
    tmp = path.with_suffix(".tmp.docx")
    with zipfile.ZipFile(path) as src, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as dst:
        for item in src.infolist():
            data = src.read(item.filename)
            if item.filename == "word/document.xml":
                text = data.decode("utf-8")
                text = text.replace("<wp:inline", "<wp:anchor")
                text = text.replace("</wp:inline>", "</wp:anchor>")
                data = text.encode("utf-8")
            dst.writestr(item, data)
    shutil.move(str(tmp), str(path))


def make_pptx(
    path: Path,
    *,
    title: str = "",
    language: str = "",
    slides: int = 1,
    slide_titles: bool = True,
    pictures: int = 0,
    picture_alt: str | None = None,
) -> Path:
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    if title:
        prs.core_properties.title = title
    if language:
        prs.core_properties.language = language
    layout = prs.slide_layouts[5]  # "Title Only"
    for index in range(slides):
        slide = prs.slides.add_slide(layout)
        if slide_titles and slide.shapes.title is not None:
            slide.shapes.title.text = f"Slide {index + 1} Title"
        for _ in range(pictures):
            pic = slide.shapes.add_picture(BytesIO(TINY_PNG), Emu(0), Emu(0))
            c_nv_pr = pic._element.nvPicPr.cNvPr
            if picture_alt is not None:
                c_nv_pr.set("descr", picture_alt)
            else:
                # python-pptx's add_picture writes a non-empty default descr
                # (e.g. the image filename); strip it so picture_alt=None
                # actually yields no alt text.
                c_nv_pr.attrib.pop("descr", None)
                c_nv_pr.attrib.pop("title", None)
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return path


def make_xlsx(
    path: Path,
    *,
    title: str = "",
    language: str = "",
    data_rows: int = 3,
    data_cols: int = 2,
    header_behaviors: bool = True,
) -> Path:
    from openpyxl import Workbook

    wb = Workbook()
    if title:
        wb.properties.title = title
    if language:
        wb.properties.language = language
    ws = wb.active
    for row in range(1, data_rows + 1):
        for col in range(1, data_cols + 1):
            ws.cell(row=row, column=col, value=f"h{col}" if row == 1 else f"v{row}.{col}")
    if header_behaviors and data_rows > 1:
        ws.freeze_panes = "A2"
        ws.auto_filter.ref = ws.dimensions
        ws.print_title_rows = "1:1"
    path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(path))
    return path


def make_fake_ole2(path: Path) -> Path:
    """A legacy .doc-shaped byte blob for FR8 guard tests."""
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_bytes(_LEGACY_OLE2_MAGIC + b"\x00" * 64)
    return path
